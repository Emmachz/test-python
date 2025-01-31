from pptx import Presentation
import sys
import os

def extract_unique_links_from_pptx(pptx_path):
    if not os.path.exists(pptx_path):
        print("Le fichier spécifié n'existe pas.")
        return []

    prs = Presentation(pptx_path)
    links = []  # Liste pour conserver l'ordre des liens extraits
    seen_links = set()  # Ensemble pour éviter les doublons

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.hyperlink and run.hyperlink.address:
                            link = run.hyperlink.address
                            if link not in seen_links:
                                links.append(link)  # Ajout à la liste si non présent
                                seen_links.add(link)  # Ajout à l'ensemble pour suivi

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.hyperlink and run.hyperlink.address:
                                    link = run.hyperlink.address
                                    if link not in seen_links:
                                        links.append(link)  # Ajout à la liste si non présent
                                        seen_links.add(link)  # Ajout à l'ensemble pour suivi

    return links  # Liste ordonnée sans doublons

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("Usage: python extract_links.py <fichier.pptx>")
        sys.exit(1)

    pptx_file = sys.argv[1]
    extracted_links = extract_unique_links_from_pptx(pptx_file)

    if extracted_links:
        print("Liens uniques extraits du fichier :", pptx_file)
        for link in extracted_links:
            print(link)
    else:
        print("Aucun lien trouvé dans le fichier.")
